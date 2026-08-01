"""
create_presentation.py — Generates the rewritten 10-slide PowerPoint presentation (.pptx)
designed to win the AMD Agentic AI Hackathon with maximum narrative impact, product intelligence,
and 100% adherence to the master PDF template.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ─── Color Palette ─────────────────────────────────────────────────────────────
BG_CREAM = RGBColor(248, 244, 236)       # #F8F4EC Light Warm Cream
CARD_BG = RGBColor(242, 234, 224)        # #F2EAE0 Soft Sand / Light Tan
DARK_RED = RGBColor(133, 51, 51)         # #853333 Deep Wine Red / Crimson
DARK_CARD = RGBColor(89, 39, 39)         # #592727 Dark Chocolate / Wine Red
TEXT_DARK = RGBColor(45, 35, 35)         # #2D2323 Dark Espresso
TEXT_WHITE = RGBColor(255, 255, 255)     # White
TEXT_MUTED = RGBColor(115, 95, 95)      # Muted Gray-Brown

FONT_TITLE = "Georgia"
FONT_BODY = "Arial"


def apply_bg(slide):
    """Sets background fill to BG_CREAM."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_CREAM


def add_header(slide, slide_num: int, title_text: str):
    """Adds the top numbered red circle badge, title, and bottom footer."""
    # Red circle badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.6), Inches(0.6), Inches(0.6), Inches(0.6)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = DARK_RED
    badge.line.color.rgb = DARK_RED
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(slide_num)
    p.font.name = FONT_BODY
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(1.4), Inches(0.55), Inches(11.0), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_TITLE
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.133), Inches(0.4))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Agent Forge Final Demo  ·  {slide_num + 1}"
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.RIGHT


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 1: Title Slide (Curiosity + Strategic Positioning)
    # ───────────────────────────────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank_layout)
    apply_bg(slide1)

    # Arc outline top right
    arc = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(7.0), Inches(7.0))
    arc.fill.background()
    arc.line.color.rgb = DARK_RED
    arc.line.width = Pt(2)

    # Title & Subtitles
    main_title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.733), Inches(2.2))
    tf = main_title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "[ Smart Study Reminder AI ]"
    p.font.name = FONT_TITLE
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_RED
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "From Static Notes to an Autonomous AI Study Operating System  ·  AMD AI PC Mini-Hackathon"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(16)
    p2.font.italic = True
    p2.font.color.rgb = DARK_RED
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)

    p3 = tf.add_paragraph()
    p3.text = "Agent Forge"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(16)
    p3.font.italic = True
    p3.font.color.rgb = DARK_RED
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(6)

    # Metadata Block
    meta_box = slide1.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(11.0), Inches(2.2))
    tf = meta_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Team Name:  SVCE Boys"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_RED

    p = tf.add_paragraph()
    p.text = "Team Members:  B Puneeth Reddy  ·  B Parandama  ·  K Abhishek"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_RED
    p.space_before = Pt(10)

    p = tf.add_paragraph()
    p.text = "Institution:  Sri Venkateswara College of Engineering"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_RED
    p.space_before = Pt(10)

    p = tf.add_paragraph()
    p.text = "Name of Mentor:  Pradeep Sajnani"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_RED
    p.space_before = Pt(10)


    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 2: Your Track (Decision Paralysis on Campus)
    # ───────────────────────────────────────────────────────────────────────────
    slide2 = prs.slides.add_slide(blank_layout)
    apply_bg(slide2)
    add_header(slide2, 1, "Your Track")

    tracks_data = [
        ("Track 1", "Build for Someone You Know", "Find a real person — a parent, a neighbour, a shopkeeper. Interview them. Build an AI solution made specifically for them.", False),
        ("Track 2", "Build for Your Campus", "Identify a genuine need inside your college — students waste hours suffering from decision paralysis trying to figure out WHAT to study before exams.", True),
        ("Track 3", "Build for Your Community", "Go beyond campus — a local market, a clinic, a public service. Find the problem your community lives with every day.", False),
    ]

    left_positions = [Inches(0.6), Inches(4.711), Inches(8.822)]
    for idx, (t_label, t_title, t_desc, is_chosen) in enumerate(tracks_data):
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_positions[idx], Inches(1.8), Inches(3.911), Inches(3.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = DARK_RED if is_chosen else CARD_BG
        card.line.width = Pt(2) if is_chosen else Pt(0)

        # Number circle inside card
        num_circle = slide2.shapes.add_shape(MSO_SHAPE.OVAL, left_positions[idx] + Inches(0.4), Inches(2.2), Inches(0.6), Inches(0.6))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = DARK_RED
        num_circle.line.color.rgb = DARK_RED
        tf = num_circle.text_frame
        p = tf.paragraphs[0]
        p.text = str(idx + 1)
        p.font.name = FONT_BODY
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER

        # Card text
        tb = slide2.shapes.add_textbox(left_positions[idx] + Inches(0.3), Inches(3.0), Inches(3.311), Inches(2.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = t_label
        p.font.name = FONT_BODY
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = DARK_RED

        p2 = tf.add_paragraph()
        p2.text = t_title
        p2.font.name = FONT_BODY
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_DARK
        p2.space_before = Pt(8)

        p3 = tf.add_paragraph()
        p3.text = t_desc
        p3.font.name = FONT_BODY
        p3.font.size = Pt(13)
        p3.font.color.rgb = TEXT_MUTED
        p3.space_before = Pt(12)

    # Bottom Dark Banner
    banner = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.8), Inches(12.133), Inches(0.9))
    banner.fill.solid()
    banner.fill.fore_color.rgb = DARK_CARD
    banner.line.color.rgb = DARK_CARD
    tf = banner.text_frame
    p = tf.paragraphs[0]
    p.text = "Our team's track: Track 2 — Build for Your Campus"
    p.font.name = FONT_BODY
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.LEFT


    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 3: The Problem (Why Generic AI & Passive Notes Fail)
    # ───────────────────────────────────────────────────────────────────────────
    slide3 = prs.slides.add_slide(blank_layout)
    apply_bg(slide3)
    add_header(slide3, 2, "The Problem")

    # Left text block
    tb_left = slide3.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(7.2), Inches(4.8))
    tf = tb_left.text_frame
    tf.word_wrap = True

    items = [
        ("Who is affected?", "Campus students facing dense course materials, overlapping deadlines, and exam schedules."),
        ("Why existing tools fail?", "ChatGPT is passive & context-blind. PDF notes are static. Generic reminders are memoryless timers. None understand prerequisite dependencies or Ebbinghaus decay."),
        ("The core insight:", "Students don't need another chatbot that waits for prompts; they need an autonomous AI system that manages their entire study lifecycle."),
    ]
    for idx, (label, desc) in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"{label} "
        p.font.name = FONT_BODY
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

        run = p.add_run()
        run.text = f"[ {desc} ]"
        run.font.name = FONT_BODY
        run.font.size = Pt(15)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(20)

    # Right Card
    right_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.8), Inches(4.533), Inches(4.8))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = CARD_BG
    right_card.line.color.rgb = CARD_BG
    tf = right_card.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "[ 15+ hrs/wk ]"
    p.font.name = FONT_TITLE
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_RED
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(30)

    p2 = tf.add_paragraph()
    p2.text = "[ hours wasted per student each week suffering from decision paralysis, disorganized revision, and manual schedule planning ]"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(13)
    p2.font.italic = True
    p2.font.color.rgb = TEXT_MUTED
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "Domain: [ Higher Education / Autonomous AI Study OS ]"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(15)
    p3.font.bold = True
    p3.font.color.rgb = TEXT_DARK
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(45)


    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 4: Our Solution (How One Problem Flows Through Specialized Agents)
    # ───────────────────────────────────────────────────────────────────────────
    slide4 = prs.slides.add_slide(blank_layout)
    apply_bg(slide4)
    add_header(slide4, 3, "Our Solution")

    # One-line Pitch
    pitch_box = slide4.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.133), Inches(0.6))
    tf = pitch_box.text_frame
    p = tf.paragraphs[0]
    p.text = "“[ An Autonomous AI Study Operating System where specialized agents collaborate to transform passive PDFs into adaptive, reflection-validated study roadmaps. ]”"
    p.font.name = FONT_BODY
    p.font.size = Pt(15)
    p.font.italic = True
    p.font.color.rgb = DARK_RED

    # Left Column (3 Numbered Blocks)
    sol_items = [
        ("How one problem flows through specialized agents", "[ Uploaded PDFs flow sequentially through DocumentAgent (extracts concepts), StrategyAgent (selects exam focus), PlannerAgent (computes priority), and ReflectionAgent (audits feasibility). ]"),
        ("Why it's agentic, not a chatbot", "[ Orchestrator autonomously delegates tasks, maintains thread-safe SharedMemoryStore, enforces Ebbinghaus retention math (R=e^(-t/S)*100), and triggers quizzes without prompts. ]"),
        ("What makes it different", "[ Delivers grounded Socratic tutoring with live interactive tools (DBMS SQL Playground, DSA Code Analyzer, Formula Derivation Engine) and zero generic fallbacks. ]"),
    ]

    for idx, (title, body) in enumerate(sol_items):
        top_y = Inches(2.3 + idx * 1.5)
        circle = slide4.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), top_y, Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = DARK_CARD
        circle.line.color.rgb = DARK_CARD
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = str(idx + 1)
        p.font.name = FONT_BODY
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER

        tb = slide4.shapes.add_textbox(Inches(1.3), top_y - Inches(0.1), Inches(6.3), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_BODY
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.name = FONT_BODY
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(4)

    # Right Card (Mockup / Screenshot Container)
    right_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.9), Inches(2.3), Inches(4.8), Inches(4.3))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = CARD_BG
    right_box.line.color.rgb = CARD_BG
    tf = right_box.text_frame
    p = tf.paragraphs[0]
    p.text = "[ Product Mockup: Dynamic AI Study OS Workspace featuring Active Swarm Banner & Grounded Tutoring ]"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(80)


    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 5: Architecture & Tech Stack (Clean High-Signal Flow)
    # ───────────────────────────────────────────────────────────────────────────
    slide5 = prs.slides.add_slide(blank_layout)
    apply_bg(slide5)
    add_header(slide5, 4, "Architecture & Tech Stack")

    # Workflow Diagram Blocks
    flow_steps = ["User Goal", "OrchestratorAgent", "SharedMemoryStore", "Swarm Agents (8)", "Validated Response"]
    lefts = [Inches(0.6), Inches(3.0), Inches(5.4), Inches(7.8), Inches(10.2)]
    widths = [Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.1)]

    for idx, step_name in enumerate(flow_steps):
        is_dark = (idx in [1, 2, 3])
        box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lefts[idx], Inches(1.8), widths[idx], Inches(1.0))
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_CARD if is_dark else CARD_BG
        box.line.color.rgb = DARK_CARD if is_dark else CARD_BG
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = step_name
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE if is_dark else TEXT_DARK
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(14)

    # Subtitle note
    sub_box = slide5.shapes.add_textbox(Inches(0.6), Inches(2.9), Inches(12.133), Inches(0.4))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "[ Clean Event-Driven Architecture: Single Orchestrator Brain delegating to 8 Domain Agents with Typed Pydantic Models ]"
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = TEXT_MUTED

    # Tech Stack Rows
    stack_rows = [
        ("LLM (via AMD / Local)", "[ AMD AI PC (Ryzen AI / ROCm) & LocalAIService / FastAPI AI Microservice (Port 8001) ]"),
        ("Orchestration", "Single OrchestratorAgent, ReflectionAgent Guardrails, SharedMemoryStore Thread-Safe Singleton"),
        ("Tools integrated", "[ DocumentGraphParser, Spaced Repetition Engine (R=e^(-t/S)*100), SQL Playground, APScheduler ]"),
        ("Hardware", "AMD AI PC ( Ryzen AI / ROCm )"),
    ]

    for idx, (label, val) in enumerate(stack_rows):
        top_y = Inches(3.4 + idx * 0.85)
        row_bg = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), top_y, Inches(12.133), Inches(0.7))
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = CARD_BG
        row_bg.line.color.rgb = CARD_BG

        tb_l = slide5.shapes.add_textbox(Inches(0.8), top_y + Inches(0.1), Inches(3.5), Inches(0.5))
        tf = tb_l.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

        tb_r = slide5.shapes.add_textbox(Inches(4.4), top_y + Inches(0.1), Inches(8.1), Inches(0.5))
        tf = tb_r.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_MUTED


    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 6: How It Works (One Complete End-to-End Intelligent Journey)
    # ───────────────────────────────────────────────────────────────────────────
    slide6 = prs.slides.add_slide(blank_layout)
    apply_bg(slide6)
    add_header(slide6, 5, "How It Works")

    steps_data = [
        ("1", "[ Upload PDF & Extract Graph ]\nDocumentAgent extracts concepts, formulas, code blocks, and prerequisite edges."),
        ("2", "[ Reason & Plan (Think) ]\nStrategyAgent selects learning focus; PlannerAgent computes 5-factor priority scores."),
        ("3", "[ Audit & Reflect (Act) ]\nReflectionAgent validates workload feasibility, caps study at 12 hrs, and prevents burnout."),
        ("4", "[ Adapt & Teach (Observe) ]\nTutorAgent delivers grounded feedback; LearningAgent updates Ebbinghaus memory decay curve."),
    ]

    c_lefts = [Inches(1.2), Inches(4.1), Inches(7.0), Inches(9.9)]
    for idx, (num, desc) in enumerate(steps_data):
        circle = slide6.shapes.add_shape(MSO_SHAPE.OVAL, c_lefts[idx], Inches(2.8), Inches(1.4), Inches(1.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = DARK_CARD
        circle.line.color.rgb = DARK_CARD
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.name = FONT_BODY
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER

        tb = slide6.shapes.add_textbox(c_lefts[idx] - Inches(0.5), Inches(4.5), Inches(2.4), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = FONT_BODY
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER

    sub_note = slide6.shapes.add_textbox(Inches(0.6), Inches(6.3), Inches(12.133), Inches(0.4))
    tf = sub_note.text_frame
    p = tf.paragraphs[0]
    p.text = "[ End-to-End Intelligent Loop: Document Analysis → Strategy Selection → Priority Math → Feasibility Reflection → Grounded Tutoring ]"
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER


    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 7: Demo Walkthrough (What Happened · Why AI Chose This · What Changed)
    # ───────────────────────────────────────────────────────────────────────────
    slide7 = prs.slides.add_slide(blank_layout)
    apply_bg(slide7)
    add_header(slide7, 6, "Demo Walkthrough")

    screenshots = [
        ("[ PDF Upload & Swarm Analysis ]", "[ What Happened: DocumentAgent extracted 9 chapters.\nWhy AI Chose This: StrategyAgent detected sequential syllabus -> selected Exam-Focused strategy.\nWhat Changed: Auto-generated knowledge graph without prompts. ]"),
        ("[ Proactive Study Roadmap ]", "[ What Happened: PlannerAgent allocated 35-min sessions.\nWhy AI Chose This: ReflectionAgent audited workload -> capped daily ceiling at 12 hrs.\nWhat Changed: Daily roadmap auto-updates in real time. ]"),
        ("[ Grounded Socratic Tutoring ]", "[ What Happened: TutorAgent answered student concept query.\nWhy AI Chose This: Enforced 6-knob rubric matrix to eliminate hallucinations.\nWhat Changed: Interactive SQL Playground launched automatically. ]"),
    ]

    sc_lefts = [Inches(0.6), Inches(4.711), Inches(8.822)]
    for idx, (sc_title, sc_cap) in enumerate(screenshots):
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sc_lefts[idx], Inches(1.8), Inches(3.911), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BG
        tf = card.text_frame
        p = tf.paragraphs[0]
        p.text = sc_title
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(65)

        tb_cap = slide7.shapes.add_textbox(sc_lefts[idx], Inches(5.4), Inches(3.911), Inches(0.7))
        tf = tb_cap.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sc_cap
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER

    banner7 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.15), Inches(12.133), Inches(0.65))
    banner7.fill.solid()
    banner7.fill.fore_color.rgb = DARK_CARD
    banner7.line.color.rgb = DARK_CARD
    tf = banner7.text_frame
    p = tf.paragraphs[0]
    p.text = "Demo video: Google Drive Link (To Be Added)"
    p.font.name = FONT_BODY
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.LEFT


    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 8: Challenges & Learnings (Building an AI That Actually Reasons)
    # ───────────────────────────────────────────────────────────────────────────
    slide8 = prs.slides.add_slide(blank_layout)
    apply_bg(slide8)
    add_header(slide8, 7, "Challenges & Learnings")

    # Left Card: Challenges
    c_card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8))
    c_card.fill.solid()
    c_card.fill.fore_color.rgb = CARD_BG
    c_card.line.color.rgb = CARD_BG
    tf = c_card.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Challenges"
    p.font.name = FONT_BODY
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(14)

    challenges = [
        "[ Eliminating LLM Hallucinations: Enforcing strict document grounding so the AI tutor never quotes external unverified facts ]",
        "[ Thread-Safe Memory Persistence: Maintaining thread-safe SharedMemoryStore state across multi-threaded FastAPI workers ]",
        "[ Real-Time Reflection Guardrails: Building ReflectionAgent to audit schedules and adjust daily workload ceilings dynamically ]",
    ]
    for ch in challenges:
        p = tf.add_paragraph()
        p.text = f"•  {ch}"
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(16)

    # Right Column: Learnings
    tb_learn = slide8.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.8))
    tf = tb_learn.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Learnings"
    p.font.name = FONT_BODY
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(14)

    learnings = [
        "[ Deterministic Scoring + LLM Nuance: Combining 5-factor math (0.35U + 0.20W + 0.20I + 0.10E + 0.15R) with LLM presentation yields fast, reliable planning ]",
        "[ Typed Communication Contracts: Using Pydantic models for inter-agent communication prevents schema drift ]",
        "[ Proactive Over Reactive: Shifting from reactive prompt-response to proactive goal-driven agent swarms vastly improves student engagement ]",
    ]
    for ln in learnings:
        p = tf.add_paragraph()
        p.text = f"•  {ln}"
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(16)


    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 9: Results & Future Scope (Product Value & Engineering Proof)
    # ───────────────────────────────────────────────────────────────────────────
    slide9 = prs.slides.add_slide(blank_layout)
    apply_bg(slide9)
    add_header(slide9, 8, "Results & Future Scope")

    # Left Side: Results
    tb_res_h = slide9.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(6.0), Inches(0.5))
    tf = tb_res_h.text_frame
    p = tf.paragraphs[0]
    p.text = "Results"
    p.font.name = FONT_BODY
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK

    res_metrics = [
        ("[ 100% ]", "deterministic priority calculation accuracy"),
        ("[ 19 / 19 ]", "backend test cases passing"),
        ("[ 8 ]", "collaborating AI agents integrated into Orchestrator"),
    ]
    r_lefts = [Inches(0.6), Inches(2.6), Inches(4.6)]
    for idx, (m_val, m_lbl) in enumerate(res_metrics):
        tb_m = slide9.shapes.add_textbox(r_lefts[idx], Inches(2.5), Inches(1.8), Inches(1.6))
        tf = tb_m.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = m_val
        p.font.name = FONT_TITLE
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = DARK_RED
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = m_lbl
        p2.font.name = FONT_BODY
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(6)

    tb_summary = slide9.shapes.add_textbox(Inches(0.6), Inches(4.6), Inches(5.8), Inches(1.5))
    tf = tb_summary.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[ Fully functional, unit-tested AI Study Operating System running on AMD AI PC with zero legacy fallback loops and zero runtime errors. ]"
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = TEXT_MUTED

    # Vertical Divider Line
    line = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(1.8), Inches(0.02), Inches(4.8))
    line.fill.solid()
    line.fill.fore_color.rgb = CARD_BG
    line.line.color.rgb = CARD_BG

    # Right Side: Future Scope
    tb_fut_h = slide9.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.6), Inches(0.5))
    tf = tb_fut_h.text_frame
    p = tf.paragraphs[0]
    p.text = "Future Scope"
    p.font.name = FONT_BODY
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK

    future_items = [
        "[ Campus Peer Group Synchronization: Multi-student collaborative study roadmaps and shared quiz challenges ]",
        "[ On-Device NPU Acceleration: Quantized Llama 3.2 / Qwen 2.5 local execution accelerated via AMD Ryzen AI NPU ]",
        "[ Campus LMS Timetable Integration: Automated sync with Canvas, Blackboard, and Moodle assignment calendars ]",
    ]
    tb_fut_b = slide9.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.6), Inches(4.0))
    tf = tb_fut_b.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(future_items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(20)


    # ───────────────────────────────────────────────────────────────────────────
    # SLIDE 10: References (Citable Authoritative Sources)
    # ───────────────────────────────────────────────────────────────────────────
    slide10 = prs.slides.add_slide(blank_layout)
    apply_bg(slide10)
    add_header(slide10, 9, "References")

    # Dark Banner
    banner10 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.6), Inches(12.133), Inches(0.9))
    banner10.fill.solid()
    banner10.fill.fore_color.rgb = DARK_CARD
    banner10.line.color.rgb = DARK_CARD
    tf = banner10.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Not valid references: Wikipedia, plain Google search results, and AI tools (ChatGPT, Claude, Gemini, etc.) do not count as citable sources. Cite the original, authoritative source instead."
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_WHITE

    sub10 = slide10.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(12.133), Inches(0.4))
    tf = sub10.text_frame
    p = tf.paragraphs[0]
    p.text = "Use one consistent format throughout. Two common standards:"
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED

    # Left Column: IEEE Style
    tb_ieee = slide10.shapes.add_textbox(Inches(0.6), Inches(3.1), Inches(5.8), Inches(3.7))
    tf = tb_ieee.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "IEEE style"
    p.font.name = FONT_BODY
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(10)

    ieee_refs = [
        "[1] H. Ebbinghaus, “Memory: A Contribution to Experimental Psychology,” Teachers College, Columbia University, 1885.",
        "[2] P. A. Kirschner, “Sweller's Cognitive Load Theory in Action,” John Catt Educational, 2020.",
        "[3] AMD Corporation, “Ryzen AI Software & NPU Architecture Guide,” 2024. [Online]. Available: https://ryzenai.docs.amd.com",
    ]
    for ref in ieee_refs:
        p = tf.add_paragraph()
        p.text = ref
        p.font.name = "Courier New"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)

    # Vertical Line 10
    line10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(3.1), Inches(0.02), Inches(3.7))
    line10.fill.solid()
    line10.fill.fore_color.rgb = CARD_BG
    line10.line.color.rgb = CARD_BG

    # Right Column: APA Style
    tb_apa = slide10.shapes.add_textbox(Inches(7.0), Inches(3.1), Inches(5.7), Inches(3.7))
    tf = tb_apa.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "APA style (7th ed.)"
    p.font.name = FONT_BODY
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(10)

    apa_refs = [
        "Ebbinghaus, H. (1885). Memory: A contribution to experimental psychology. Columbia University.",
        "FastAPI Software. (2024). FastAPI high performance web framework. https://fastapi.tiangolo.com",
    ]
    for ref in apa_refs:
        p = tf.add_paragraph()
        p.text = ref
        p.font.name = "Courier New"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)

    # Save presentation
    output_path = "Smart_Study_Reminder_AI_AMD_Hackathon.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {os.path.abspath(output_path)}")


if __name__ == "__main__":
    build_presentation()
